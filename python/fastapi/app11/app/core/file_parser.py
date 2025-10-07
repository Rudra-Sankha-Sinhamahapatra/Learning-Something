import os
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

def extract_text_from_file(file, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()

    if ext in [".txt", ".md"]:
        return file.read().decode("utf-8")
    
    elif ext == ".pdf":
        reader = PdfReader(file)
        return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
    
    elif ext == ".docx":
        doc = DocxDocument(file)
        return "\n".join(p.text for p in doc.paragraphs)
    
    elif ext == ".pptx":
        prs = Presentation(file)
        slides = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)
        return "\n".join(slides)
    
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    

def extract_text_from_url(url: str) -> str:
    resp = requests.get(url)
    soup = BeautifulSoup(resp.text, "html.parser")
    return soup.get_text()